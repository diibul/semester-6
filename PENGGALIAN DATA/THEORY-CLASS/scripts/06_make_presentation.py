"""
06_make_presentation.py
=======================
Generate a comprehensive PPTX presentation covering all analysis phases,
including EDA, RFM, Advanced Analytics, Baseline Model, and GMM Clustering.

Author: Muhammad Iqbal Fadel
Date: May 2026 (updated June 2026)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
EDA_DIR = os.path.join(ROOT, 'outputs', 'figures', 'eda')
GMM_DIR = os.path.join(ROOT, 'outputs', 'figures', 'gmm')
MODEL_DIR = os.path.join(ROOT, 'outputs', 'figures', 'models')
ANALYSIS_DIR = os.path.join(ROOT, 'outputs', 'figures', 'analysis')
OUT = os.path.join(ROOT, 'reports', 'presentation_progres_01-06-2026.pptx')


def add_title_slide(prs, title, subtitle):
    """Add a title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide


def add_section_slide(prs, section_title, section_subtitle=""):
    """Add a section header slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[2])  # Section Header layout
    slide.shapes.title.text = section_title
    if section_subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = section_subtitle
    return slide


def add_image_slide(prs, title_text, image_path, subtitle_text=""):
    """Add a slide with an image and title."""
    if not os.path.exists(image_path):
        print(f'  [SKIP] Not found: {image_path}')
        return None

    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank with title
    slide.shapes.title.text = title_text

    # Add image centered
    left = Inches(0.4)
    top = Inches(1.4)
    max_width = Inches(9.2)
    max_height = Inches(5.5)

    try:
        from PIL import Image
        im = Image.open(image_path)
        iw, ih = im.size
        aspect = iw / ih
        if aspect > (max_width / max_height):
            width = max_width
            height = width / aspect
        else:
            height = max_height
            width = height * aspect
        # Center horizontally
        left = Inches((10 - width / Emu(914400)) / 2)
    except Exception:
        width = max_width
        height = max_height

    slide.shapes.add_picture(image_path, left, top, width=width, height=height)
    print(f'  [OK] Added slide: {title_text}')
    return slide


def main():
    print('=' * 60)
    print('  GENERATING PPTX PRESENTATION')
    print('=' * 60)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Title Slide ──────────────────────────────────────────────
    add_title_slide(
        prs,
        "Data Mining — E-commerce Customer Segmentation",
        "Analisis Transaksi Penjualan & GMM Clustering\n"
        "Muhammad Iqbal Fadel | Semester 6 | Juni 2026"
    )

    # ══════════════════════════════════════════════════════════════
    # SECTION 1: EDA
    # ══════════════════════════════════════════════════════════════
    add_section_slide(prs, "Phase 1: Exploratory Data Analysis", "12 visualisasi dari 522,601 transaksi")

    eda_slides = [
        ("Monthly Sales Trend", 'monthly_sales_trend.png'),
        ("Top 20 Products by Revenue", 'top20_products_by_revenue.png'),
        ("Top 15 Countries by Sales", 'top15_countries_by_sales.png'),
        ("Quantity Distribution", 'quantity_distribution.png'),
        ("Total Amount Boxplot (Outlier Detection)", 'totalamount_boxplot.png'),
    ]
    for title, img in eda_slides:
        add_image_slide(prs, title, os.path.join(EDA_DIR, img))

    # ══════════════════════════════════════════════════════════════
    # SECTION 2: RFM Segmentation
    # ══════════════════════════════════════════════════════════════
    add_section_slide(prs, "Phase 2: RFM Customer Segmentation", "5 segmen pelanggan berdasarkan Recency, Frequency, Monetary")

    rfm_slides = [
        ("RFM Segment Distribution", 'rfm_segment_counts.png'),
        ("RFM Segment Averages Heatmap", 'rfm_segment_averages_heatmap.png'),
    ]
    for title, img in rfm_slides:
        add_image_slide(prs, title, os.path.join(EDA_DIR, img))

    # ══════════════════════════════════════════════════════════════
    # SECTION 3: Advanced Analytics
    # ══════════════════════════════════════════════════════════════
    add_section_slide(prs, "Phase 3: Advanced Analytics", "Basket Analysis, Cohort Retention, Sales Forecasting")

    advanced_slides = [
        ("Top Product Pairs (Basket Analysis)", 'basket_top_pairs.png'),
        ("Customer Retention Cohort Heatmap", 'cohort_retention_heatmap.png'),
        ("Monthly Sales Forecast (Linear Trend)", 'monthly_sales_forecast.png'),
    ]
    for title, img in advanced_slides:
        add_image_slide(prs, title, os.path.join(EDA_DIR, img))

    # ══════════════════════════════════════════════════════════════
    # SECTION 4: Baseline ML Model
    # ══════════════════════════════════════════════════════════════
    add_section_slide(prs, "Phase 4: Machine Learning — Baseline RandomForest",
                      "Cross-validation, Feature Importance, ROC AUC")

    model_slides = [
        ("ROC Curve — Baseline RandomForest", 'baseline_roc.png'),
        ("Feature Importance (Gini)", 'feature_importance.png'),
        ("Confusion Matrix", 'confusion_matrix.png'),
        ("5-Fold Cross-Validation Scores", 'cross_validation_scores.png'),
    ]
    for title, img in model_slides:
        add_image_slide(prs, title, os.path.join(MODEL_DIR, img))

    # ══════════════════════════════════════════════════════════════
    # SECTION 5: Feature Selection
    # ══════════════════════════════════════════════════════════════
    analysis_slides = [
        ("Feature Selection Comparison", 'feature_selection_comparison.png'),
        ("Model Performance Comparison", 'model_performance_comparison.png'),
    ]
    has_analysis = any(os.path.exists(os.path.join(ANALYSIS_DIR, img)) for _, img in analysis_slides)
    if has_analysis:
        add_section_slide(prs, "Phase 5: Feature Selection Analysis", "Correlation-based, Backward Elimination, Model-based")
        for title, img in analysis_slides:
            add_image_slide(prs, title, os.path.join(ANALYSIS_DIR, img))

    # ══════════════════════════════════════════════════════════════
    # SECTION 6: GMM CLUSTERING (HIGHLIGHT)
    # ══════════════════════════════════════════════════════════════
    add_section_slide(prs, "HIGHLIGHT: GMM Clustering",
                      "Gaussian Mixture Model — Probabilistic Customer Segmentation\n"
                      "9 cluster optimal via BIC/AIC | 7 visualisasi profesional")

    gmm_slides = [
        ("GMM: BIC & AIC Cluster Selection", 'gmm_bic_aic_selection.png'),
        ("GMM: PCA 2D Projection", 'gmm_cluster_scatter_2d.png'),
        ("GMM: PCA 3D Projection", 'gmm_cluster_scatter_3d.png'),
        ("GMM: Cluster Profiles — Radar Chart", 'gmm_cluster_profiles_radar.png'),
        ("GMM: Cluster Distribution", 'gmm_cluster_distribution.png'),
        ("GMM: Membership Probability Heatmap", 'gmm_probability_heatmap.png'),
        ("GMM vs RFM — Cross-tabulation Comparison", 'gmm_vs_rfm_comparison.png'),
    ]
    for title, img in gmm_slides:
        add_image_slide(prs, title, os.path.join(GMM_DIR, img))

    # ── Save ──────────────────────────────────────────────────────
    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    prs.save(OUT)
    slide_count = len(prs.slides)
    print(f'\n  [OK] Saved PPTX: {OUT}')
    print(f'  [OK] Total slides: {slide_count}')


if __name__ == '__main__':
    main()
